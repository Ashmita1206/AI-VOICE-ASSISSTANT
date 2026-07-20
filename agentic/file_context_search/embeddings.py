"""
Text Extractor + Embedding Generator
=====================================

Provides two responsibilities that naturally live together:

1. ``extract_text(path)``  — pulls readable text from any supported file type.
2. ``generate_embedding(text)`` — encodes text into a float32 numpy vector
   using a local sentence-transformers model (no external API).

The SentenceTransformer model is loaded *once* as a module-level singleton
to avoid repeated disk I/O across calls.
"""

from __future__ import annotations

import logging
import os
import re
from pathlib import Path
from typing import Optional

import numpy as np

logger = logging.getLogger(__name__)

# ── Configuration ──────────────────────────────────────────────────────────
EMBEDDING_MODEL_NAME: str = "all-MiniLM-L6-v2"
MAX_EXTRACT_CHARS: int = 8000    # characters to read from each file
MAX_EMBED_CHARS: int = 512       # chars fed to the embedding model


# ── Lazy-loaded SentenceTransformer singleton ──────────────────────────────
_model = None
_model_load_attempted: bool = False


def _get_model():
    """Return the SentenceTransformer model, loading it on first call."""
    global _model, _model_load_attempted
    if _model is not None:
        return _model
    if _model_load_attempted:
        return None  # failed before; don't retry on every call
    _model_load_attempted = True
    try:
        from sentence_transformers import SentenceTransformer
        logger.info("[EMBEDDINGS] Loading model: %s", EMBEDDING_MODEL_NAME)
        _model = SentenceTransformer(EMBEDDING_MODEL_NAME)
        logger.info("[EMBEDDINGS] Model loaded successfully.")
        return _model
    except ImportError:
        logger.warning(
            "[EMBEDDINGS] sentence-transformers not installed. "
            "Run: pip install sentence-transformers"
        )
        return None
    except Exception as exc:
        logger.error("[EMBEDDINGS] Failed to load model: %s", exc)
        return None


# ── Text Extraction ────────────────────────────────────────────────────────

def _extract_pdf(path: str) -> str:
    """Extract text from a PDF using PyMuPDF (fitz)."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(path)
        parts: list[str] = []
        char_count = 0
        for page in doc:
            text = page.get_text("text")
            parts.append(text)
            char_count += len(text)
            if char_count >= MAX_EXTRACT_CHARS:
                break
        doc.close()
        return "\n".join(parts)[:MAX_EXTRACT_CHARS]
    except ImportError:
        logger.debug("PyMuPDF not installed; skipping PDF: %s", path)
        return ""
    except Exception as exc:
        logger.debug("PDF extraction failed for %s: %s", path, exc)
        return ""


def _extract_docx(path: str) -> str:
    """Extract text from a DOCX using python-docx."""
    try:
        import docx
        doc = docx.Document(path)
        parts: list[str] = []
        char_count = 0
        for para in doc.paragraphs:
            parts.append(para.text)
            char_count += len(para.text)
            if char_count >= MAX_EXTRACT_CHARS:
                break
        return "\n".join(parts)[:MAX_EXTRACT_CHARS]
    except ImportError:
        logger.debug("python-docx not installed; skipping DOCX: %s", path)
        return ""
    except Exception as exc:
        logger.debug("DOCX extraction failed for %s: %s", path, exc)
        return ""


def _extract_pptx(path: str) -> str:
    """Extract text from a PPTX using python-pptx."""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts: list[str] = []
        char_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    parts.append(shape.text)
                    char_count += len(shape.text)
                    if char_count >= MAX_EXTRACT_CHARS:
                        break
            if char_count >= MAX_EXTRACT_CHARS:
                break
        return "\n".join(parts)[:MAX_EXTRACT_CHARS]
    except ImportError:
        logger.debug("python-pptx not installed; skipping PPTX: %s", path)
        return ""
    except Exception as exc:
        logger.debug("PPTX extraction failed for %s: %s", path, exc)
        return ""


def _extract_xlsx(path: str) -> str:
    """Extract text from an XLSX/XLS using openpyxl."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts: list[str] = []
        char_count = 0
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_str = " ".join(str(c) for c in row if c is not None)
                parts.append(row_str)
                char_count += len(row_str)
                if char_count >= MAX_EXTRACT_CHARS:
                    break
            if char_count >= MAX_EXTRACT_CHARS:
                break
        wb.close()
        return "\n".join(parts)[:MAX_EXTRACT_CHARS]
    except ImportError:
        logger.debug("openpyxl not installed; skipping XLSX: %s", path)
        return ""
    except Exception as exc:
        logger.debug("XLSX extraction failed for %s: %s", path, exc)
        return ""


def _extract_plain(path: str) -> str:
    """Read plain-text files (txt, md, csv, json, py, html, …)."""
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as fh:
            return fh.read(MAX_EXTRACT_CHARS)
    except Exception as exc:
        logger.debug("Plain-text extraction failed for %s: %s", path, exc)
        return ""


def extract_text(path: str) -> str:
    """Dispatch to the appropriate extractor based on file extension.

    Parameters
    ----------
    path:
        Absolute path to the file.

    Returns
    -------
    str
        Extracted text (up to ``MAX_EXTRACT_CHARS`` characters).
        Returns empty string on failure — never raises.
    """
    ext = Path(path).suffix.lower().lstrip(".")

    if ext == "pdf":
        return _extract_pdf(path)
    elif ext in {"docx", "doc"}:
        return _extract_docx(path)
    elif ext in {"pptx", "ppt"}:
        return _extract_pptx(path)
    elif ext in {"xlsx", "xls"}:
        return _extract_xlsx(path)
    else:
        # txt, md, csv, json, py, js, html, xml, yaml, log, etc.
        return _extract_plain(path)


# ── Keyword Extraction ─────────────────────────────────────────────────────

_STOP_WORDS: frozenset[str] = frozenset({
    "the", "a", "an", "and", "or", "but", "is", "are", "was", "were",
    "be", "been", "being", "have", "has", "had", "do", "does", "did",
    "will", "would", "shall", "should", "may", "might", "must", "can",
    "could", "of", "in", "on", "at", "to", "for", "with", "by", "from",
    "up", "about", "into", "through", "during", "before", "after",
    "above", "below", "as", "if", "then", "that", "this", "these",
    "those", "it", "its", "i", "you", "he", "she", "we", "they",
    "me", "him", "her", "us", "them", "my", "your", "his", "our",
    "their", "not", "no", "so", "than", "just", "also", "very",
    "more", "most", "such", "there", "here", "when", "where", "who",
    "which", "what", "how", "all", "each", "any", "some", "new",
})


def extract_keywords(text: str, max_keywords: int = 50) -> str:
    """Extract meaningful keywords from text.

    Returns a space-joined string of the top ``max_keywords`` tokens.
    """
    tokens = re.findall(r"[a-zA-Z]{3,}", text.lower())
    filtered = [t for t in tokens if t not in _STOP_WORDS]
    # Frequency-rank
    from collections import Counter
    freq = Counter(filtered)
    top = [word for word, _ in freq.most_common(max_keywords)]
    return " ".join(top)


def extract_summary(text: str, max_chars: int = 200) -> str:
    """Return the first meaningful sentence(s) of extracted text."""
    if not text:
        return ""
    # Strip excessive whitespace
    clean = re.sub(r"\s+", " ", text).strip()
    # Take the first max_chars characters; try to end at a sentence boundary
    if len(clean) <= max_chars:
        return clean
    truncated = clean[:max_chars]
    last_period = max(truncated.rfind("."), truncated.rfind("!"), truncated.rfind("?"))
    if last_period > max_chars // 2:
        return truncated[:last_period + 1]
    return truncated + "…"


# ── Embedding Generation ───────────────────────────────────────────────────

def generate_embedding(text: str) -> Optional[bytes]:
    """Encode text as a float32 numpy vector and return as raw bytes.

    Parameters
    ----------
    text:
        Text to encode. Only the first ``MAX_EMBED_CHARS`` characters are used.

    Returns
    -------
    bytes or None
        Raw bytes of the float32 numpy array, suitable for storing in SQLite
        as a BLOB.  Returns ``None`` if the model is unavailable.
    """
    model = _get_model()
    if model is None:
        return None

    snippet = text[:MAX_EMBED_CHARS].strip()
    if not snippet:
        return None

    try:
        vector: np.ndarray = model.encode(snippet, convert_to_numpy=True, normalize_embeddings=True)
        return vector.astype(np.float32).tobytes()
    except Exception as exc:
        logger.warning("[EMBEDDINGS] Failed to encode text: %s", exc)
        return None


def blob_to_vector(blob: bytes) -> Optional[np.ndarray]:
    """Convert a stored BLOB back to a float32 numpy array."""
    if not blob:
        return None
    try:
        return np.frombuffer(blob, dtype=np.float32)
    except Exception:
        return None


def cosine_similarity(a: np.ndarray, b: np.ndarray) -> float:
    """Compute cosine similarity between two unit-normalised vectors."""
    try:
        dot = float(np.dot(a, b))
        norm_a = float(np.linalg.norm(a))
        norm_b = float(np.linalg.norm(b))
        if norm_a == 0.0 or norm_b == 0.0:
            return 0.0
        return dot / (norm_a * norm_b)
    except Exception:
        return 0.0
