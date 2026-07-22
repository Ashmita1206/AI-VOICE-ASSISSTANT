"""
Document Retrieval Utilities
============================

Shared utilities for hashing, text extraction, and summarization.
"""

import hashlib
import logging
import os
import re
from pathlib import Path
from typing import Optional

from agentic.document_retrieval import config

logger = logging.getLogger(__name__)

# ── Hashing ────────────────────────────────────────────────────────────────

def get_content_hash(path: str) -> str:
    """Compute SHA-256 hash of a file's contents."""
    h = hashlib.sha256()
    try:
        with open(path, "rb") as f:
            for chunk in iter(lambda: f.read(8192), b""):
                h.update(chunk)
        return h.hexdigest()
    except Exception as e:
        logger.debug(f"Failed to hash {path}: {e}")
        return ""

# ── Intelligent Text Extraction ────────────────────────────────────────────

def _extract_pdf(path: str) -> str:
    try:
        import fitz
        doc = fitz.open(path)
        parts = []
        
        # 1. Try to get Table of Contents / Bookmarks
        try:
            toc = doc.get_toc()
            if toc:
                toc_lines = [f"Heading: {item[1]}" for item in toc if len(item) > 1]
                parts.append("Table of Contents:\n" + "\n".join(toc_lines[:30]))
        except Exception:
            pass

        # 2. Extract first 5 pages, middle pages, and last 2 pages if large
        total_pages = len(doc)
        pages_to_read = []
        if total_pages <= 10:
            pages_to_read = list(range(total_pages))
        else:
            pages_to_read = list(range(min(5, total_pages)))
            mid = total_pages // 2
            pages_to_read.extend([mid - 1, mid, mid + 1])
            pages_to_read.extend(range(max(0, total_pages - 2), total_pages))

        char_count = 0
        for p_idx in sorted(set(pages_to_read)):
            if p_idx < total_pages:
                text = doc[p_idx].get_text("text")
                if text.strip():
                    parts.append(f"--- Page {p_idx + 1} ---\n{text}")
                    char_count += len(text)
                    if char_count >= config.MAX_EXTRACT_CHARS:
                        break
        doc.close()
        return "\n".join(parts)[:config.MAX_EXTRACT_CHARS]
    except Exception as e:
        logger.debug(f"PDF extraction failed for {path}: {e}")
        return ""

def _extract_docx(path: str) -> str:
    try:
        import docx
        doc = docx.Document(path)
        parts = []
        char_count = 0

        # Paragraphs & Headings
        for para in doc.paragraphs:
            if not para.text.strip():
                continue
            if para.style and "heading" in para.style.name.lower():
                parts.append(f"Section Heading: {para.text}")
            else:
                parts.append(para.text)
            char_count += len(para.text)
            if char_count >= config.MAX_EXTRACT_CHARS:
                break

        # Tables
        if char_count < config.MAX_EXTRACT_CHARS:
            for table in doc.tables:
                for row in table.rows:
                    row_txt = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_txt:
                        parts.append(row_txt)
                        char_count += len(row_txt)
                        if char_count >= config.MAX_EXTRACT_CHARS:
                            break

        return "\n".join(parts)[:config.MAX_EXTRACT_CHARS]
    except Exception as e:
        logger.debug(f"DOCX extraction failed for {path}: {e}")
        return ""

def _extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        char_count = 0
        for idx, slide in enumerate(prs.slides, start=1):
            slide_parts = [f"--- Slide {idx} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())
            # Speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_parts.append(f"Notes: {notes}")
            
            slide_txt = "\n".join(slide_parts)
            parts.append(slide_txt)
            char_count += len(slide_txt)
            if char_count >= config.MAX_EXTRACT_CHARS:
                break
        return "\n".join(parts)[:config.MAX_EXTRACT_CHARS]
    except Exception as e:
        logger.debug(f"PPTX extraction failed for {path}: {e}")
        return ""

def _extract_xlsx(path: str) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts = []
        char_count = 0
        for sheet in wb.worksheets:
            parts.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                row_str = " ".join(str(c) for c in row if c is not None)
                if row_str.strip():
                    parts.append(row_str)
                    char_count += len(row_str)
                    if char_count >= config.MAX_EXTRACT_CHARS:
                        break
            if char_count >= config.MAX_EXTRACT_CHARS:
                break
        wb.close()
        return "\n".join(parts)[:config.MAX_EXTRACT_CHARS]
    except Exception as e:
        logger.debug(f"XLSX extraction failed for {path}: {e}")
        return ""

def _extract_csv(path: str) -> str:
    try:
        import csv
        parts = []
        char_count = 0
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            for idx, row in enumerate(reader):
                if idx > 200:
                    break
                row_str = " | ".join(row)
                if row_str.strip():
                    parts.append(row_str)
                    char_count += len(row_str)
                    if char_count >= config.MAX_EXTRACT_CHARS:
                        break
        return "\n".join(parts)[:config.MAX_EXTRACT_CHARS]
    except Exception as e:
        logger.debug(f"CSV extraction failed for {path}: {e}")
        return ""

def _extract_ipynb(path: str) -> str:
    try:
        import json
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            nb = json.load(f)
        parts = []
        char_count = 0
        for cell in nb.get("cells", []):
            source = "".join(cell.get("source", []))
            if source.strip():
                cell_type = cell.get("cell_type", "code")
                parts.append(f"[{cell_type.upper()}]\n{source}")
                char_count += len(source)
                if char_count >= config.MAX_EXTRACT_CHARS:
                    break
        return "\n".join(parts)[:config.MAX_EXTRACT_CHARS]
    except Exception as e:
        logger.debug(f"IPYNB extraction failed for {path}: {e}")
        return ""

def _extract_plain(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return f.read(config.MAX_EXTRACT_CHARS)
    except Exception as e:
        logger.debug(f"Plain text extraction failed for {path}: {e}")
        return ""

def extract_text(path: str) -> str:
    """Extract text intelligently based on file type."""
    ext = Path(path).suffix.lower().lstrip(".")
    if ext == "pdf":
        return _extract_pdf(path)
    elif ext in ("doc", "docx"):
        return _extract_docx(path)
    elif ext in ("ppt", "pptx"):
        return _extract_pptx(path)
    elif ext in ("xls", "xlsx"):
        return _extract_xlsx(path)
    elif ext == "csv":
        return _extract_csv(path)
    elif ext == "ipynb":
        return _extract_ipynb(path)
    else:
        return _extract_plain(path)

# ── Summarization & Keyword Extraction ─────────────────────────────────────

def generate_summary(text: str, max_chars: int = 250) -> str:
    """Generate a clean summary of the text, stripping HTML markup and noise."""
    if not text:
        return ""
    
    # Strip HTML tags, doctypes, scripts, and CSS
    clean = re.sub(r"<script.*?>.*?</script>", " ", text, flags=re.DOTALL | re.IGNORECASE)
    clean = re.sub(r"<style.*?>.*?</style>", " ", clean, flags=re.DOTALL | re.IGNORECASE)
    clean = re.sub(r"<[^>]+>", " ", clean)
    clean = re.sub(r"<!DOCTYPE.*?>", " ", clean, flags=re.IGNORECASE)
    clean = re.sub(r"\s+", " ", clean).strip()
    
    if len(clean) <= max_chars:
        return clean
        
    truncated = clean[:max_chars]
    # Try to end at a sentence boundary
    last_punct = max(truncated.rfind("."), truncated.rfind("!"), truncated.rfind("?"))
    if last_punct > max_chars // 2:
        return truncated[:last_punct + 1]
    return truncated + "..."

_STOP_WORDS = frozenset({
    "the", "a", "an", "and", "or", "but", "is", "are", "was", "were",
    "be", "been", "being", "have", "has", "had", "do", "does", "did",
    "will", "would", "shall", "should", "may", "might", "must", "can",
    "could", "of", "in", "on", "at", "to", "for", "with", "by", "from",
    "up", "about", "into", "through", "during", "before", "after",
    "this", "that", "it", "they", "we", "he", "she", "which", "what", "who", "where", "when", "why", "how",
    "not", "no", "yes", "so", "than", "just", "also", "very", "more", "most", "some", "any", "all", "each",
    "there", "here", "their", "our", "my", "your", "his", "her", "its",
})

def extract_keywords(text: str, top_n: int = 50) -> str:
    """Extract top keywords based on frequency."""
    if not text:
        return ""
        
    tokens = re.findall(r"\b[a-zA-Z]{3,}\b", text.lower())
    filtered = [t for t in tokens if t not in _STOP_WORDS]
    
    from collections import Counter
    freq = Counter(filtered)
    return " ".join([word for word, count in freq.most_common(top_n)])
